#!/usr/bin/env python3
"""生成 AI 聊天系统架构 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

IMG_DIR = os.path.dirname(os.path.abspath(__file__))

# 颜色
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BLUE = RGBColor(0x3b, 0x82, 0xf6)
LIGHT_GRAY = RGBColor(0x94, 0xa3, 0xb8)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AI 智能聊天系统"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "架构设计与技术实践"
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "策略模式 · RAG 知识库 · 对话记忆 · 工具调度 · 可观测性 · 动态路由"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_GRAY
    
    return slide

def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox

def add_image(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        if width:
            return slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            return slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
        else:
            return slide.shapes.add_picture(img_path, Inches(left), Inches(top))

def add_bullet_list(slide, left, top, width, height, items, size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)

# ==================== Slide 1: 封面 ====================
add_title_slide()

# ==================== Slide 2: 目录 ====================
slide = add_content_slide("目录")
add_bullet_list(slide, 1, 1.5, 11, 5, [
    "1. 系统概述与背景",
    "2. 整体架构图及说明",
    "3. 业务用例图及说明",
    "4. 协作图（时序图）及说明",
    "5. 类图及说明",
    "6. 设计模式与核心组件",
    "7. 技术栈与基础设施",
    "8. 总结与展望",
], size=20)

# ==================== Slide 3: 系统概述 ====================
slide = add_content_slide("1. 系统概述与背景", "从「调接口」到「智能体编排平台」的演进")
add_bullet_list(slide, 0.8, 1.8, 12, 5, [
    "Spring Boot 3 + React 18 + Milvus + Redis + MySQL(RDS) + RabbitMQ",
    "",
    "演进历程：",
    "  阶段一：设计模式重构（策略+工厂+外观）——消除 740 行重复代码",
    "  阶段二：RAG 知识库——让 AI 先查资料再回答",
    "  阶段三：对话记忆持久化——短期 Redis + 长期 Milvus",
    "  阶段四：体验增强——重新生成/停止/摘要/多模态",
    "  阶段五：工具调度层——让 AI 能「执行」任务",
    "  阶段六：可观测性——链路追踪/错误聚合/AI 自愈",
    "  阶段七：动态模型路由——按任务类型自动选最优模型",
    "",
    "核心指标：35 个新增文件 | 15+ 项功能 | 代码减少 37% | 8 个可独立开关的模块",
], size=15)

# ==================== Slide 4: 整体架构图 ====================
slide = add_content_slide("2. 整体架构图", "七层 Harness 智能体编排架构")
add_image(slide, os.path.join(IMG_DIR, "A_professional_software_archit_2026-08-07T20-19-51.png"), 1.5, 1.5, width=10)

# ==================== Slide 5: 架构说明 ====================
slide = add_content_slide("2.1 架构层说明", "各层职责与设计理念")
add_bullet_list(slide, 0.8, 1.5, 12, 5.5, [
    "业务层：ChatProcessor / DebateProcessor / TreeHoleService，只关心业务逻辑",
    "Harness 编排层：RAGService(知识库) + MemoryService(记忆) + ToolDispatcher(工具) + SelfHealing(自愈)",
    "动态路由层：TaskClassifier 分类任务 → ModelRouter 评分选最优模型",
    "可观测性层：TraceContext 链路追踪 + ErrorAggregator 错误聚合 + 自动自愈",
    "外观层(Facade)：LLMInvoker 统一入口，收敛 baseUrl/apiKey/统计/日志",
    "策略层(Strategy)：OpenAICompatStrategy + DoubaoStrategy，隔离供应商差异",
    "基础设施：Milvus(向量库) + Redis(缓存) + MySQL RDS(业务) + RabbitMQ(消息)",
    "",
    "设计理念：所有 AI 能力在编排层统一管理，业务层零侵入接入",
], size=15)

# ==================== Slide 6: 业务用例图 ====================
slide = add_content_slide("3. 业务用例图", "用户与管理员的功能边界")
add_image(slide, os.path.join(IMG_DIR, "A_UML_use_case_diagram_for_an__2026-08-07T20-20-13.png"), 1.5, 1.5, width=10)

# ==================== Slide 7: 用例说明 ====================
slide = add_content_slide("3.1 用例说明", "核心功能模块")
add_bullet_list(slide, 0.8, 1.5, 5.8, 5.5, [
    "用户用例：",
    "  注册/登录",
    "  AI 伙伴群聊（多 AI 角色）",
    "  个人对话空间（私密+记忆）",
    "  情绪树洞（情感倾听+RAG）",
    "  观点辩论场（正反方+共识）",
    "  AI 多人游戏（3款）",
    "  知识脉络图",
    "  图片/视频/3D 生成",
    "  历史搜索（分页弹窗）",
    "  语音输入/朗读",
    "  重新生成/停止",
], size=14)
add_bullet_list(slide, 6.8, 1.5, 5.8, 5.5, [
    "管理员用例：",
    "  模型管理（增删改查）",
    "  知识库管理（上传文档）",
    "  SQL 执行器",
    "  监控看板",
    "  调用链路追踪",
    "  错误聚合分析",
    "",
    "技术特性：",
    "  RAG 知识库增强",
    "  对话记忆持久化",
    "  工具调度（天气/计算器）",
    "  动态模型路由",
    "  AI 错误自愈",
], size=14)

# ==================== Slide 8: 协作图 ====================
slide = add_content_slide("4. 协作图（RAG 流程）", "检索增强生成的完整调用链路")
add_image(slide, os.path.join(IMG_DIR, "A_UML_collaboration_sequence_d_2026-08-07T20-21-02.png"), 1.5, 1.5, width=10)

# ==================== Slide 9: 协作图说明 ====================
slide = add_content_slide("4.1 协作流程说明", "RAG + 记忆 + 工具调用的编排流程")
add_bullet_list(slide, 0.8, 1.5, 12, 5.5, [
    "RAG 流程：",
    "  用户提问 → EmbeddingService 向量化 → Milvus 检索 top5 → 拼入 system prompt → LLMInvoker 调用 → 回答",
    "",
    "记忆流程：",
    "  保存：对话存 Redis(短期5轮) + Milvus(长期永久)",
    "  读取：Redis 取最近5轮 + Milvus 检索相关3条 → 拼入 system prompt",
    "",
    "工具调度流程：",
    "  用户提问 → LLM 判断是否需要工具 → 执行工具(天气/计算器/时间/知识库) → 结果回填 → LLM 生成最终回答",
    "",
    "自愈流程：",
    "  LLM 调用失败 → ErrorAggregator 记录 → SelfHealingService 根据错误类型换模型/降温度/重试 → 最多重试2次",
], size=15)

# ==================== Slide 10: 类图 ====================
slide = add_content_slide("5. 类图", "设计模式与核心类关系")
add_image(slide, os.path.join(IMG_DIR, "A_UML_class_diagram_for_an_AI__2026-08-07T20-20-39.png"), 1.5, 1.5, width=10)

# ==================== Slide 11: 类图说明 ====================
slide = add_content_slide("5.1 类图说明", "策略+工厂+外观三种设计模式")
add_bullet_list(slide, 0.8, 1.5, 12, 5.5, [
    "策略模式（Strategy）：",
    "  LLMStrategy 接口：invoke() / invokeStream() / supportsStream()",
    "  OpenAICompatStrategy：覆盖千问/DeepSeek/豆包，走 /chat/completions",
    "  DoubaoStrategy：走 /responses 接口，不支持真流式时伪流式降级",
    "",
    "工厂模式（Factory）：",
    "  LLMStrategyFactory：按 provider 路由策略，默认走 OpenAICompatStrategy",
    "  新增 getStrategyForTask()：VISION 任务强制用 OpenAICompatStrategy",
    "",
    "外观模式（Facade）：",
    "  LLMInvoker：统一入口，收敛 baseUrl 解析 + API Key 选择 + 调用统计 + 链路追踪 + 错误自愈",
    "  业务层只需调 llmInvoker.invoke()，不关心策略选择和横切逻辑",
], size=15)

# ==================== Slide 12: 设计模式收益 ====================
slide = add_content_slide("6. 设计模式与核心组件", "量化收益")
add_bullet_list(slide, 0.8, 1.5, 5.8, 5.5, [
    "代码收益：",
    "  重复代码消除：~740 行",
    "  业务层代码减少：37%",
    "  LLM 调用入口：4→1",
    "  新增供应商：改4文件→加1策略",
    "",
    "新增模块（35文件）：",
    "  strategy/ 策略模式(3)",
    "  factory/ 工厂模式(1)",
    "  rag/ 知识库+记忆(13)",
    "  agent/ 工具调度(7)",
    "  observability/ 可观测性(6)",
    "  router/ 动态路由(4)",
    "  util/ 工具类(1)",
], size=14)
add_bullet_list(slide, 6.8, 1.5, 5.8, 5.5, [
    "配置开关（8个模块独立控制）：",
    "  app.rag.enabled RAG知识库",
    "  app.agent.enabled 工具调度",
    "  app.observability.enabled 可观测性",
    "  app.observability.auto-heal AI自愈",
    "  app.router.enabled 动态路由",
    "  app.classifier.enabled 任务分类",
    "",
    "核心能力：",
    "  AI 记忆：10分钟→永久",
    "  知识库：无→PDF/Word/TXT",
    "  工具调用：无→4种工具",
    "  错误自愈：无→5种策略",
    "  模型选择：手动→自动路由",
], size=14)

# ==================== Slide 13: 技术栈 ====================
slide = add_content_slide("7. 技术栈与基础设施", "生产级技术选型")
add_bullet_list(slide, 0.8, 1.5, 5.8, 5.5, [
    "后端：",
    "  Spring Boot 3 + MyBatis",
    "  WebSocket (STOMP) 实时通信",
    "  Spring Security + JWT 认证",
    "  RabbitMQ 跨节点消息广播",
    "  Redis 缓存/Session/统计",
    "  MySQL (阿里云 RDS)",
    "  Milvus 2.3.4 向量数据库",
    "",
    "前端：",
    "  React 18 + Vite",
    "  STOMP.js WebSocket",
    "  Web Speech API 语音",
    "  SpeechSynthesis 朗读",
], size=14)
add_bullet_list(slide, 6.8, 1.5, 5.8, 5.5, [
    "基础设施：",
    "  应用服务器：2核1.6G (双实例)",
    "  Milvus 服务器：8G 独立部署",
    "  阿里云 RDS：VPC 内网",
    "  Docker Compose：Milvus+etcd+MinIO",
    "  Nginx 反向代理 + 负载均衡",
    "",
    "AI 能力：",
    "  千问 (qwen-plus/max/vl)",
    "  DeepSeek (对话/代码)",
    "  豆包 (doubao-seed-character)",
    "  text-embedding-v3 (向量化)",
    "  阿里云 OSS (媒体存储)",
    "  阿里云内容安全",
], size=14)

# ==================== Slide 14: 总结 ====================
slide = add_content_slide("8. 总结与展望", "从调接口到智能体的完整演进")
add_bullet_list(slide, 0.8, 1.5, 12, 3, [
    "已完成：",
    "  设计模式重构 → RAG 知识库 → 对话记忆 → 体验增强 → 工具调度 → 可观测性 → 动态路由",
    "  形成完整的 Harness 智能体编排架构，8 个模块可独立开关，零侵入接入",
    "",
    "核心价值：",
    "  模型是通用的，但数据、业务逻辑、用户体验是独特的",
    "  Harness 工程把独特性编织进 AI 能力，让通用模型变成「懂业务的专属 AI」",
], size=15)

add_bullet_list(slide, 0.8, 5, 12, 2.5, [
    "后续演进方向：",
    "  1. HTTPS 配置（解锁 iOS 语音输入）",
    "  2. Agent 链式调用（多步推理）",
    "  3. 微调专属模型（LoRA + vLLM）",
    "  4. 对话分析看板（词云/情绪/满意度）",
    "  5. 熔断降级（Resilience4j）",
    "  6. Token 计费与成本控制",
], size=14, )

# 保存
output = os.path.join(IMG_DIR, "AI聊天系统架构设计.pptx")
prs.save(output)
print(f"PPT 已生成：{output}")
